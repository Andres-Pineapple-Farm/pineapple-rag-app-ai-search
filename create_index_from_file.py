import os
from azure.identity import DefaultAzureCredential
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient
from langchain_openai import AzureOpenAIEmbeddings
from config import get_logger
import tempfile
from pathlib import Path

# Import Office document libraries
try:
    import docx  # For Word documents
    import pptx  # For PowerPoint presentations
    HAS_OFFICE_SUPPORT = True
except ImportError:
    HAS_OFFICE_SUPPORT = False
    logger = get_logger(__name__)
    logger.warning("python-docx and/or python-pptx not installed. Word and PowerPoint support disabled.")
    logger.warning("To enable, install with: pip install python-docx python-pptx")

# initialize logging object
logger = get_logger(__name__)

# Initialize embeddings
try:
    embeddings = AzureOpenAIEmbeddings(
        deployment="text-embedding-ada-002",
        model="text-embedding-ada-002",
        api_key=os.getenv("AZURE_OPENAI_API_KEY"),
        azure_endpoint=os.getenv("AZURE_OPENAI_API_BASE"),
        api_version=os.getenv("AZURE_OPENAI_API_VERSION")
    )
    # Test that embed_query method exists and works
    test_embedding = embeddings.embed_query("test")
    logger.info(f"AzureOpenAIEmbeddings initialized successfully. Vector dimension: {len(test_embedding)}")
except Exception as e:
    logger.error(f"Error initializing embeddings: {e}")
    raise

# Azure Cognitive Search configuration
search_service_endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
search_api_key = os.getenv("AZURE_SEARCH_API_KEY")

# Initialize SearchIndexClient
index_client = SearchIndexClient(endpoint=search_service_endpoint, 
                               credential=AzureKeyCredential(search_api_key))

import pandas as pd
from azure.search.documents.indexes.models import (
    SemanticSearch,
    SearchField,
    SimpleField,
    SearchableField,
    SearchFieldDataType,
    SemanticConfiguration,
    SemanticPrioritizedFields,
    SemanticField,
    VectorSearch,
    HnswAlgorithmConfiguration,
    VectorSearchAlgorithmKind,
    HnswParameters,
    VectorSearchAlgorithmMetric,
    ExhaustiveKnnAlgorithmConfiguration,
    ExhaustiveKnnParameters,
    VectorSearchProfile,
    SearchIndex,
)

def create_index_definition(index_name: str, model: str) -> SearchIndex:
    dimensions = 1536  # text-embedding-ada-002
    if model == "text-embedding-3-large":
        dimensions = 3072

    # The fields we want to index. The "embedding" field is a vector field that will
    # be used for vector search.
    fields = [
        SimpleField(name="id", type=SearchFieldDataType.String, key=True),
        SearchableField(name="content", type=SearchFieldDataType.String),
        SimpleField(name="filepath", type=SearchFieldDataType.String),
        SearchableField(name="title", type=SearchFieldDataType.String),
        SimpleField(name="url", type=SearchFieldDataType.String),
        SearchField(
            name="contentVector",
            type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
            searchable=True,
            # Size of the vector created by the text-embedding-ada-002 model.
            vector_search_dimensions=dimensions,
            vector_search_profile_name="myHnswProfile",
        ),
    ]

    # The "content" field should be prioritized for semantic ranking.
    semantic_config = SemanticConfiguration(
        name="default",
        prioritized_fields=SemanticPrioritizedFields(
            title_field=SemanticField(field_name="title"),
            keywords_fields=[],
            content_fields=[SemanticField(field_name="content")],
        ),
    )

    # For vector search, we want to use the HNSW (Hierarchical Navigable Small World)
    # algorithm (a type of approximate nearest neighbor search algorithm) with cosine
    # distance.
    vector_search = VectorSearch(
        algorithms=[
            HnswAlgorithmConfiguration(
                name="myHnsw",
                kind=VectorSearchAlgorithmKind.HNSW,
                parameters=HnswParameters(
                    m=4,
                    ef_construction=1000,
                    ef_search=1000,
                    metric=VectorSearchAlgorithmMetric.COSINE,
                ),
            ),
            ExhaustiveKnnAlgorithmConfiguration(
                name="myExhaustiveKnn",
                kind=VectorSearchAlgorithmKind.EXHAUSTIVE_KNN,
                parameters=ExhaustiveKnnParameters(metric=VectorSearchAlgorithmMetric.COSINE),
            ),
        ],
        profiles=[
            VectorSearchProfile(
                name="myHnswProfile",
                algorithm_configuration_name="myHnsw",
            ),
            VectorSearchProfile(
                name="myExhaustiveKnnProfile",
                algorithm_configuration_name="myExhaustiveKnn",
            ),
        ],
    )

    # Create the semantic settings with the configuration
    semantic_search = SemanticSearch(configurations=[semantic_config])

    # Create the search index definition
    return SearchIndex(
        name=index_name,
        fields=fields,
        semantic_search=semantic_search,
        vector_search=vector_search,
    )

# define a function for indexing a markdown file, that chunks the content
# and generates vector embeddings for each chunk
def create_docs_from_markdown(path: str, model: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[dict[str, any]]:
    # Read the markdown file
    try:
        # Check file extension to ensure we're only trying to read text files
        file_ext = Path(path).suffix.lower()
        binary_extensions = ['.docx', '.pptx', '.pdf']
        
        if file_ext in binary_extensions:
            raise ValueError(f"Cannot read binary file format {file_ext} directly with create_docs_from_markdown. Use the appropriate processor instead.")
            
        with open(path, 'r', encoding='utf-8') as file:
            markdown_content = file.read()
        
        logger.info(f"Successfully read markdown file: {path} ({len(markdown_content)} chars)")
    except UnicodeDecodeError as e:
        error_msg = f"File {path} appears to be binary, not text. Use appropriate processor for this file type."
        logger.error(error_msg)
        raise ValueError(error_msg) from e
    except Exception as e:
        logger.error(f"Error reading markdown file: {e}")
        raise    # Check if the content is from a PowerPoint presentation (looking for slide markers)
    if "## Slide " in markdown_content:
        # For PowerPoint content, split by slides
        slide_pattern = "## Slide "
        pages = markdown_content.split(slide_pattern)
        
        # The first element might be the header without a slide number
        if not pages[0].strip().startswith("## Slide"):
            header = pages[0]
            pages = pages[1:]
        else:
            header = ""
        
        # Set a flag to indicate this is slide content, not page content
        is_slide_content = True
        content_type = "slides"
    elif "## Page " in markdown_content:
        # Regular page-based content (PDFs)
        page_pattern = "## Page "
        pages = markdown_content.split(page_pattern)
        
        # The first element might be the header without a page number
        if not pages[0].strip().startswith("## Page"):
            header = pages[0]
            pages = pages[1:]
        else:
            header = ""
            
        is_slide_content = False
        content_type = "pages"
    else:
        # Content without any page/slide markers (like Word documents)
        # Process the entire content as a single document or chunk it by sections
        logger.info("No page or slide markers found, treating as a single document with sections")
        
        # For Word docs and other content without markers, split by heading patterns
        # This is a simple approach - we look for markdown headings (starting with #)
        heading_pattern = "\n#"  # New line followed by #
        sections = markdown_content.split(heading_pattern)
        
        # Process the first section (might be empty or contain content before first heading)
        if sections[0].strip():
            # If there's content before the first heading, treat it as its own section
            header = sections[0]
            pages = sections[1:]
            # Add the heading marker back to each section except the first one
            pages = ["#" + page for page in pages]
        else:
            # If first section is empty, ignore it
            header = ""
            pages = sections[1:]
            # Add the heading marker back to each section
            pages = ["#" + page for page in pages]
        
        # If we couldn't find any sections, treat the whole document as one section
        if not pages:
            pages = [markdown_content]
            header = ""
            
        is_slide_content = False
        content_type = "sections"
    
    items = []
    for i, page_content in enumerate(pages):
        # Skip empty pages/slides
        if not page_content.strip():
            continue
          # Extract page/slide number
        if is_slide_content:
            # For slides, use the slide number directly
            slide_num = str(i + 1)  # +1 because slides are 1-indexed in the original format
            content = page_content
            # Check if there's a slide title (### heading) and extract it
            slide_title = None
            lines = content.split('\n')
            for line in lines:
                if line.startswith('### '):
                    slide_title = line.replace('### ', '').strip()
                    break
        else:
            # For regular pages, extract page number if possible
            try:
                if page_content.strip() and page_content.strip()[0].isdigit():
                    page_num = page_content.split('\n')[0].strip()
                    content = '\n'.join(page_content.split('\n')[1:])
                else:
                    page_num = str(i)
                    content = page_content
            except IndexError:
                page_num = str(i)
                content = page_content
        
        # Further chunk the content if it's too large
        chunks = []
        if len(content) > chunk_size:
            # Simple chunking by splitting into smaller pieces with overlap
            for j in range(0, len(content), chunk_size - chunk_overlap):
                chunk = content[j:j + chunk_size]
                if chunk.strip():  # Skip empty chunks
                    chunks.append(chunk)
        else:
            chunks = [content]
          # Create a document for each chunk
        for j, chunk in enumerate(chunks):
            try:
                # Generate a unique ID for each chunk
                if is_slide_content:
                    chunk_id = f"slide{slide_num}_chunk{j}"
                    
                    # Create a title based on slide number, title, and chunk
                    if slide_title:
                        if len(chunks) > 1:
                            title = f"Slide {slide_num}: {slide_title} - Chunk {j+1}"
                        else:
                            title = f"Slide {slide_num}: {slide_title}"
                    else:
                        if len(chunks) > 1:
                            title = f"Slide {slide_num} - Chunk {j+1}"
                        else:
                            title = f"Slide {slide_num}"
                    
                    # Set the URL to point to the specific slide
                    url = f"/document/slide-{slide_num}"
                elif content_type == "pages":
                    chunk_id = f"page{page_num}_chunk{j}"
                    
                    # Create a title based on page number and chunk
                    if len(chunks) > 1:
                        title = f"Page {page_num} - Chunk {j+1}"
                    else:
                        title = f"Page {page_num}"
                    
                    # Set the URL to point to the specific page
                    url = f"/document/page-{page_num}"
                else:
                    # For sections (like Word documents)
                    section_num = i + 1  # 1-indexed section number
                    chunk_id = f"section{section_num}_chunk{j}"
                    
                    # Try to extract a heading from the chunk
                    section_title = None
                    lines = chunk.split('\n')
                    for line in lines[:3]:  # Check first few lines for a heading
                        if line.startswith('#'):
                            # Remove heading markers and spaces
                            section_title = line.lstrip('#').strip()
                            break
                    
                    # Create a title based on section number, title (if any), and chunk
                    if section_title:
                        if len(chunks) > 1:
                            title = f"Section {section_num}: {section_title} - Chunk {j+1}"
                        else:
                            title = f"Section {section_num}: {section_title}"
                    else:
                        if len(chunks) > 1:
                            title = f"Section {section_num} - Chunk {j+1}"
                        else:
                            title = f"Section {section_num}"
                    
                    # Set the URL to point to the specific section
                    url = f"/document/section-{section_num}"
                
                # Generate embeddings for the chunk using embed_query
                embedding_vector = embeddings.embed_query(chunk)
                  # Create and add the record
                rec = {
                    "id": chunk_id,
                    "content": chunk,
                    "filepath": path,
                    "title": title,
                    "url": url,
                    "contentVector": embedding_vector,
                }
                items.append(rec)
            except Exception as e:
                logger.error(f"Error generating embeddings for chunk: {e}")
                raise
    
    logger.info(f"Created {len(items)} document chunks from markdown file")
    return items

# define a function for indexing a csv file, that adds each row as a document
# and generates vector embeddings for the specified content_column
def create_docs_from_csv(path: str, content_column: str, model: str) -> list[dict[str, any]]:
    try:
        products = pd.read_csv(path)
        logger.info(f"Successfully read CSV file: {path}")
    except Exception as e:
        logger.error(f"Error reading CSV file: {e}")
        raise
        
    items = []
    for product in products.to_dict("records"):
        try:
            content = product[content_column]
            id = str(product["id"])
            title = product["name"]
            url = f"/products/{title.lower().replace(' ', '-')}"
            
            # Generate embeddings for the content using embed_query
            embedding_vector = embeddings.embed_query(content)
            
            rec = {
                "id": id,
                "content": content,
                "filepath": f"{title.lower().replace(' ', '-')}",
                "title": title,
                "url": url,
                "contentVector": embedding_vector,
            }
            items.append(rec)
        except Exception as e:
            logger.error(f"Error processing CSV row: {e}")
            raise

    logger.info(f"Created {len(items)} documents from CSV file")
    return items

def create_index_from_file(index_name, file_path, file_type='markdown'):
    """
    Create an Azure AI Search index from a file.
    
    Args:
        index_name: The name of the search index to create
        file_path: Path to the file to process
        file_type: Type of file to process ('markdown', 'csv', 'word', 'powerpoint')
    """
    # If a search index already exists, delete it:
    try:
        index_definition = index_client.get_index(index_name)
        index_client.delete_index(index_name)
        logger.info(f"🗑️  Found existing index named '{index_name}', and deleted it")
    except Exception:
        logger.info(f"🗑️  No existing index named '{index_name}' found, so no need to delete it")

    # Create an empty search index
    try:
        index_definition = create_index_definition(index_name, model=embeddings.model)
        index_client.create_index(index_definition)
        logger.info(f"Created index '{index_name}'")
    except Exception as e:
        logger.error(f"Error creating index: {e}")
        raise
        
    # Based on file type, process the file differently
    docs = []
    try:
        if file_type.lower() == 'csv':
            # Create documents from the CSV file, generating vector embeddings
            docs = create_docs_from_csv(path=file_path, content_column="description", model=embeddings.model)
        elif file_type.lower() == 'markdown':
            # Create documents from the markdown file, chunking and generating vector embeddings
            docs = create_docs_from_markdown(path=file_path, model=embeddings.model)
        elif file_type.lower() == 'word':
            # Create documents from the Word document, generating vector embeddings
            docs = create_docs_from_word(path=file_path, model=embeddings.model)
        elif file_type.lower() == 'powerpoint':
            # Create documents from the PowerPoint document, generating vector embeddings
            docs = create_docs_from_powerpoint(path=file_path, model=embeddings.model)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        raise
        
    if not docs or len(docs) == 0:
        logger.warning(f"No documents were created from file: {file_path}")
        return

    # Log document structure for debugging
    if docs:
        sample_doc = docs[0]
        logger.info(f"Sample document keys: {list(sample_doc.keys())}")
        logger.info(f"Document count: {len(docs)}")

    # Add the documents to the index using the Azure AI Search client
    try:
        search_client = SearchClient(
            endpoint=search_service_endpoint,
            index_name=index_name,
            credential=AzureKeyCredential(search_api_key)
        )

        search_client.upload_documents(docs)
        logger.info(f"➕ Uploaded {len(docs)} documents to '{index_name}' index")
    except Exception as e:
        logger.error(f"Error uploading documents to index: {e}")
        raise

# if __name__ == "__main__":
#     import argparse

#     parser = argparse.ArgumentParser()
#     parser.add_argument(
#         "--index-name",
#         type=str,
#         help="index name to use when creating the AI Search index",
#         default=os.environ["AISEARCH_INDEX_NAME"],
#     )
#     parser.add_argument(
#         "--file", type=str, help="path to data file for creating search index", required=True
#     )
#     parser.add_argument(
#         "--file-type", type=str, choices=["csv", "markdown"], default="csv",
#         help="type of file to process (csv or markdown)"
#     )

#     print("I'm in the main function")

#     args = parser.parse_args()
#     index_name = args.index_name
#     file_path = args.file
#     file_type = args.file_type
    
#     logger.info(f"Creating index from {file_type} file: {file_path}")
#     create_index_from_file(index_name, file_path, file_type)
#     logger.info("Index created.")

def create_docs_from_word(path: str, model: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[dict[str, any]]:
    """
    Process a Word document and generate vector embeddings for each chunk.
    
    This function extracts text from Word documents while preserving structure
    such as headings, paragraphs, and tables, then converts to markdown format
    for consistent processing.
    
    Args:
        path: Path to the Word document
        model: The embedding model name to use
        chunk_size: Base size for text chunks
        chunk_overlap: Overlap between consecutive chunks
        
    Returns:
        List of document records with embeddings and metadata
    """
    if not HAS_OFFICE_SUPPORT:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")
    
    logger.info(f"Processing Word document: {path}")
    
    # Load the document
    doc = docx.Document(path)
    
    # Extract text with structure preservation
    all_text = []
    
    # Process the document title if available
    title = doc.core_properties.title
    if title:
        all_text.append(f"# {title}")
    else:
        # Use filename as title if no document title is set
        all_text.append(f"# {Path(path).stem}")
    
    # Extract document headings and paragraphs with structure
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
            
        # Check if it's a heading
        if para.style.name.startswith('Heading'):
            # Get heading level from style name (e.g., 'Heading 1' -> 1)
            try:
                heading_level = int(para.style.name.split(' ')[-1])
                prefix = '#' * (heading_level + 1)  # +1 because we used H1 for the document title
            except (ValueError, IndexError):
                prefix = '##'  # Default to H2 if we can't determine the level
                
            all_text.append(f"{prefix} {para.text}")
        else:
            # Regular paragraph
            all_text.append(para.text)
    
    # Extract tables
    for table in doc.tables:
        table_rows = []
        header_row = []
        
        # Process header row
        for cell in table.rows[0].cells:
            header_row.append(cell.text.strip())
        
        if header_row:
            table_rows.append('| ' + ' | '.join(header_row) + ' |')
            # Add separator row
            table_rows.append('| ' + ' | '.join(['---' for _ in header_row]) + ' |')
        
        # Process data rows
        for row in table.rows[1:]:
            row_cells = [cell.text.strip() for cell in row.cells]
            if any(row_cells):  # Only add if row has content
                table_rows.append('| ' + ' | '.join(row_cells) + ' |')
        
        if table_rows:
            all_text.append("\n".join(table_rows))
    
    # Convert to markdown format for processing
    markdown_content = "\n\n".join(all_text)
    
    # Create a temporary markdown file
    temp_markdown_path = Path(tempfile.gettempdir()) / f"{Path(path).stem}_converted.md"
    with open(temp_markdown_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    
    logger.info(f"Created temporary markdown from Word document at: {temp_markdown_path}")
    logger.info(f"Markdown content length: {len(markdown_content)} characters")
    logger.info(f"First 200 characters of markdown: {markdown_content[:200]}")
    
    # Use optimized chunking parameters for Word docs
    word_chunk_size = min(chunk_size * 1.2, 1800)  # Larger chunks for Word docs
    word_chunk_overlap = min(chunk_overlap * 1.5, 350)  # More overlap to maintain context
    
    # Use existing markdown processing
    try:
        docs = create_docs_from_markdown(
            path=str(temp_markdown_path),
            model=model,
            chunk_size=int(word_chunk_size),
            chunk_overlap=int(word_chunk_overlap)
        )
        
        # Update filepath to point to original Word document
        for doc in docs:
            doc['filepath'] = path
            # Remove content_type as it's not in the index schema
            # doc['content_type'] = 'word'
        
        logger.info(f"Successfully created {len(docs)} document chunks from Word")
        return docs
    except Exception as e:
        logger.error(f"Error processing Word document markdown: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise

def create_docs_from_powerpoint(path: str, model: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[dict[str, any]]:
    """
    Process a PowerPoint presentation and generate vector embeddings for each chunk.
    
    This function extracts text from PowerPoint slides while preserving structure
    such as titles, content, and notes, then converts to markdown format for 
    consistent processing.
    
    Args:
        path: Path to the PowerPoint document
        model: The embedding model name to use
        chunk_size: Base size for text chunks
        chunk_overlap: Overlap between consecutive chunks
        
    Returns:
        List of document records with embeddings and metadata
    """
    if not HAS_OFFICE_SUPPORT:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    
    logger.info(f"Processing PowerPoint presentation: {path}")
    
    # Load the presentation
    presentation = pptx.Presentation(path)
    
    # Extract text with structure preservation
    all_text = []
    
    # Add the presentation title
    title = Path(path).stem
    all_text.append(f"# {title}")
    
    # Process each slide
    for i, slide in enumerate(presentation.slides):
        # Add slide number as heading
        slide_number = i + 1
        all_text.append(f"## Slide {slide_number}")
        
        # Extract slide title if available
        if slide.shapes.title and slide.shapes.title.text:
            all_text.append(f"### {slide.shapes.title.text.strip()}")
        
        # Extract text from all shapes in the slide
        texts = []
        found_content = False
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
                found_content = True
        
        # Add a placeholder if no text content was found
        if not found_content:
            texts.append("[Slide may contain images or other non-text content]")
        
        if texts:
            all_text.append("\n".join(texts))
        
        # Extract notes if available
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                all_text.append(f"**Notes:**\n\n{notes_text}")
    
    # Convert to markdown format for processing
    markdown_content = "\n\n".join(all_text)
    
    # Create a temporary markdown file
    temp_markdown_path = Path(tempfile.gettempdir()) / f"{Path(path).stem}_converted.md"
    with open(temp_markdown_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    
    logger.info(f"Created temporary markdown from PowerPoint at: {temp_markdown_path}")
    logger.info(f"Markdown content length: {len(markdown_content)} characters")
    logger.info(f"First 200 characters of markdown: {markdown_content[:200]}")
    
    # Use optimized chunking parameters for PowerPoint docs
    # PowerPoint slides are more fragmented, so we use smaller chunks with more overlap
    ppt_chunk_size = min(chunk_size * 0.8, 1500)  # Smaller chunks than Word docs
    ppt_chunk_overlap = min(chunk_overlap * 1.5, 350)  # More overlap to maintain context
    
    # Use existing markdown processing
    try:
        docs = create_docs_from_markdown(
            path=str(temp_markdown_path),
            model=model,
            chunk_size=int(ppt_chunk_size),
            chunk_overlap=int(ppt_chunk_overlap)
        )
        
        # Update filepath to point to original PowerPoint document
        for doc in docs:
            doc['filepath'] = path
            # Remove content_type as it's not in the index schema
            # doc['content_type'] = 'powerpoint'
        
        logger.info(f"Successfully created {len(docs)} document chunks from PowerPoint")
        
        if not docs:
            logger.warning(f"No document chunks were created from PowerPoint file: {path}")
        
        return docs
    except Exception as e:
        logger.error(f"Error processing PowerPoint markdown: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise
