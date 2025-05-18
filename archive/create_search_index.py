import os
from azure.identity import DefaultAzureCredential
from azure.core.credentials import AzureKeyCredential

from azure.search.documents import SearchClient
from azure.search.documents.indexes import SearchIndexClient

from azure.search.documents.indexes.models import (
    SemanticSearch,
    SearchField,
    SimpleField,
    SearchableField,
    SearchFieldDataType,
    SemanticConfiguration,
    SemanticPrioritizedFields,
    SemanticField,
    VectorSearch,
    HnswAlgorithmConfiguration,
    VectorSearchAlgorithmKind,
    HnswParameters,
    VectorSearchAlgorithmMetric,
    ExhaustiveKnnAlgorithmConfiguration,
    ExhaustiveKnnParameters,
    VectorSearchProfile,
    SearchIndex,
)
import pandas as pd
import PyPDF2  # For native PDF handling
import subprocess
import tempfile
from pathlib import Path
import importlib.util  # For dynamically importing the pdf-to-markdown module
from config import get_logger
from file_handler import detect_file_type, validate_file  # Import file detection utilities
# Additional imports for Office documents
try:
    import docx  # For Word documents
    import pptx  # For PowerPoint presentations
    HAS_OFFICE_SUPPORT = True
except ImportError:
    HAS_OFFICE_SUPPORT = False
    logger.warning("python-docx and/or python-pptx not installed. Word and PowerPoint support disabled.")
    logger.warning("To enable, install with: pip install python-docx python-pptx")

from langchain_openai import AzureOpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter


# initialize logging object
logger = get_logger(__name__)

# Initialize embeddings
embeddings = AzureOpenAIEmbeddings(
    deployment="text-embedding-ada-002",
    model="text-embedding-ada-002",
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    azure_endpoint=os.getenv("AZURE_OPENAI_API_BASE"),
    api_version=os.getenv("AZURE_OPENAI_API_VERSION")
)

# Azure Cognitive Search configuration
search_service_endpoint = os.getenv("AZURE_SEARCH_ENDPOINT")
search_api_key = os.getenv("AZURE_SEARCH_API_KEY")
index_name = "document-index"

# Initialize SearchIndexClient and SearchClient
index_client = SearchIndexClient(endpoint=search_service_endpoint, 
                                 credential=AzureKeyCredential(search_api_key))

# Define the index name and the model to use for vector embeddings
def create_index_definition(index_name: str, model: str) -> SearchIndex:
    dimensions = 1536  # text-embedding-ada-002
    if model == "text-embedding-3-large":
        dimensions = 3072

    # The fields we want to index. The "embedding" field is a vector field that will
    # be used for vector search.
    fields = [
        SimpleField(name="id", type=SearchFieldDataType.String, key=True),
        SearchableField(name="content", type=SearchFieldDataType.String),
        SimpleField(name="filepath", type=SearchFieldDataType.String),
        SearchableField(name="title", type=SearchFieldDataType.String),
        SimpleField(name="url", type=SearchFieldDataType.String),
        SearchField(
            name="contentVector",
            type=SearchFieldDataType.Collection(SearchFieldDataType.Single),
            searchable=True,
            # Size of the vector created by the text-embedding-ada-002 model.
            vector_search_dimensions=dimensions,
            vector_search_profile_name="myHnswProfile",
        ),
    ]

    # The "content" field should be prioritized for semantic ranking.
    semantic_config = SemanticConfiguration(
        name="default",
        prioritized_fields=SemanticPrioritizedFields(
            title_field=SemanticField(field_name="title"),
            keywords_fields=[],
            content_fields=[SemanticField(field_name="content")],
        ),
    )

    # For vector search, we want to use the HNSW (Hierarchical Navigable Small World)
    # algorithm (a type of approximate nearest neighbor search algorithm) with cosine
    # distance.
    vector_search = VectorSearch(
        algorithms=[
            HnswAlgorithmConfiguration(
                name="myHnsw",
                kind=VectorSearchAlgorithmKind.HNSW,
                parameters=HnswParameters(
                    m=4,
                    ef_construction=1000,
                    ef_search=1000,
                    metric=VectorSearchAlgorithmMetric.COSINE,
                ),
            ),
            ExhaustiveKnnAlgorithmConfiguration(
                name="myExhaustiveKnn",
                kind=VectorSearchAlgorithmKind.EXHAUSTIVE_KNN,
                parameters=ExhaustiveKnnParameters(metric=VectorSearchAlgorithmMetric.COSINE),
            ),
        ],
        profiles=[
            VectorSearchProfile(
                name="myHnswProfile",
                algorithm_configuration_name="myHnsw",
            ),
            VectorSearchProfile(
                name="myExhaustiveKnnProfile",
                algorithm_configuration_name="myExhaustiveKnn",
            ),
        ],
    )

    # Create the semantic settings with the configuration
    semantic_search = SemanticSearch(configurations=[semantic_config])

    # Create the search index definition
    return SearchIndex(
        name=index_name,
        fields=fields,
        semantic_search=semantic_search,
        vector_search=vector_search,
    )

# define a function for indexing a markdown file, that chunks the content
# and generates vector embeddings for each chunk
def create_docs_from_markdown(path: str, 
                              embedding_model, 
                              chunk_size: int = 1000, 
                              chunk_overlap: int = 200) -> list[dict[str, any]]:
    
    # Read the markdown file
    with open(path, 'r', encoding='utf-8') as file:
        markdown_content = file.read()
    
    # Define the header structure for our MarkdownHeaderTextSplitter
    # This handles headers from H1 to H6
    headers_to_split_on = [
        ("#", "Header 1"),
        ("##", "Header 2"),
        ("###", "Header 3"),
        ("####", "Header 4"),
        ("#####", "Header 5"),
        ("######", "Header 6"),
    ]
    
    # First, split the markdown by headers
    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
    header_splits = markdown_splitter.split_text(markdown_content)
    
    # Then use RecursiveCharacterTextSplitter for further chunking with specific separators
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n## ", "\n\n", "\n", " ", ""],
        keep_separator=True
    )
    
    items = []
    chunk_id_counter = 0
    
    # Process each header section
    for i, header_split in enumerate(header_splits):
        # Get metadata from the header split
        header_metadata = header_split.metadata
        heading_text = ""
        
        # Build a hierarchical heading path
        for header_level in ["Header 1", "Header 2", "Header 3", "Header 4", "Header 5", "Header 6"]:
            if header_level in header_metadata and header_metadata[header_level]:
                if heading_text:
                    heading_text += " > "
                heading_text += header_metadata[header_level]
        
        # If no heading found, use a default
        if not heading_text:
            heading_text = f"Section {i+1}"
        
        # Further chunk the content if needed
        content = header_split.page_content
        chunks = text_splitter.split_text(content)
        
        # Create a document for each chunk
        for j, chunk in enumerate(chunks):
            # Skip empty chunks
            if not chunk.strip():
                continue
                
            # Generate a unique ID for each chunk
            chunk_id = f"chunk_{chunk_id_counter}"
            chunk_id_counter += 1
            
            # Create a title based on header and chunk number
            if len(chunks) > 1:
                title = f"{heading_text} - Part {j+1}"
            else:
                title = heading_text
            
            # Set the URL to point to the section
            url = f"/document/{heading_text.lower().replace(' ', '-').replace('>', '-')}"
            
            # Generate embeddings for the chunk using LangChain's embedding model
            embedding_vector = embedding_model.embed_query(chunk)
            
            # Create and add the record
            rec = {
                "id": chunk_id,
                "content": chunk,
                "filepath": path,
                "title": title,
                "url": url,
                "contentVector": embedding_vector,
            }
            items.append(rec)
    
    return items

# define a function for indexing a csv file, that adds each row as a document
# and generates vector embeddings for the specified content_column
def create_docs_from_csv(path: str, content_column: str, embedding_model) -> list[dict[str, any]]:
    products = pd.read_csv(path)
    items = []
    for product in products.to_dict("records"):
        content = product[content_column]
        id = str(product["id"])
        title = product["name"]
        url = f"/products/{title.lower().replace(' ', '-')}"
        # Use LangChain's embed_query method
        embedding_vector = embedding_model.embed_query(content)
        rec = {
            "id": id,
            "content": content,
            "filepath": f"{title.lower().replace(' ', '-')}",
            "title": title,
            "url": url,
            "contentVector": embedding_vector,
        }
        items.append(rec)

    return items

def convert_pdf_to_text(pdf_path: str) -> str:
    """
    Convert a PDF file to text using PyPDF2.
    
    Args:
        pdf_path: Path to the PDF file to convert
    
    Returns:
        Extracted text content from the PDF file
    """
    text = ""
    with open(pdf_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

# define a function for indexing a pdf file, that handles both native PDFs and image-based PDFs
def create_docs_from_pdf(path: str,
                         embedding_model,
                         chunk_size: int = 1000, 
                         chunk_overlap: int = 200) -> list[dict[str, any]]:
    """
    Process a PDF file and generate vector embeddings for each chunk.
    
    This function intelligently handles two types of PDFs:
    1. Native PDFs with extractable text - processed using PyPDF2 to extract text directly
    2. Image-based PDFs (scans) that require OCR - processed using Azure Document Intelligence
    
    Detection Logic:
    - The function attempts to extract text from the first page using PyPDF2
    - If extractable text is found, it processes the file as a native PDF
    - If no text is found or extraction fails, it treats it as an image-based PDF
    
    Chunking Strategy:
    - For PDFs, larger chunks (1.5x) with more overlap (2x) are used to preserve context
    - Maximum chunk size is capped at 2000 characters
    - The function uses LangChain's MarkdownHeaderTextSplitter to preserve document structure
    
    Requirements:
    - For native PDFs: PyPDF2
    - For image PDFs: Azure Document Intelligence credentials (DOCUMENTINTELLIGENCE_ENDPOINT and DOCUMENTINTELLIGENCE_API_KEY)
    
    Args:
        path: Path to the PDF file
        embedding_model: The model to use for generating embeddings
        chunk_size: Base size for each text chunk (will be adjusted for PDFs)
        chunk_overlap: Base overlap between consecutive chunks (will be adjusted for PDFs)
        
    Returns:
        A list of document records with embeddings and metadata
    """
    logger.info(f"Processing PDF file: {path}")
    
    # Check if it's a native PDF with extractable text
    is_native_pdf = False
    
    try:
        # Try to extract text with PyPDF2
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            # Check if we can extract text from the first page
            sample_text = pdf_reader.pages[0].extract_text().strip()
            is_native_pdf = len(sample_text) > 0
            logger.info(f"PDF type detection: {'Native PDF' if is_native_pdf else 'Image-based PDF'}")
    except Exception as e:
        logger.warning(f"Error checking PDF type: {e}")
        is_native_pdf = False
    
    if is_native_pdf:
        # Process native PDF by extracting text
        logger.info("Processing native PDF with extractable text")
        
        # Extract text from all pages
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            all_text = []
            
            for i, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text.strip():  # Only add non-empty pages
                    all_text.append(f"## Page {i+1}\n{page_text}")
            
            # Join all the text with double newlines
            full_text = "\n\n".join(all_text)
        
        # Create a temporary markdown file
        temp_markdown_path = Path(tempfile.gettempdir()) / f"{Path(path).stem}_converted.md"
        with open(temp_markdown_path, 'w', encoding='utf-8') as f:
            f.write(f"# Converted from PDF To Markdown\n\n{full_text}")
        
        logger.info(f"Created temporary markdown from native PDF at: {temp_markdown_path}")
          # Process the markdown file using the existing function
        # For PDFs, we typically want slightly larger chunks to preserve context
        pdf_chunk_size = min(chunk_size * 1.5, 2000)  # Larger chunks, but max 2000 chars
        pdf_chunk_overlap = min(chunk_overlap * 2, 400)  # More overlap to maintain context
        
        docs = create_docs_from_markdown(
            path=str(temp_markdown_path),
            embedding_model=embedding_model,
            chunk_size=int(pdf_chunk_size),
            chunk_overlap=int(pdf_chunk_overlap)
        )
        
        # Update filepath to point to the original PDF
        for doc in docs:
            doc['filepath'] = path
        
        return docs
        
    else:
        # Process image-based PDF using pdftomarkdown.py
        logger.info("Processing image-based PDF with Document Intelligence")
        
        # Get Document Intelligence credentials from environment
        endpoint = os.getenv("DOCUMENTINTELLIGENCE_ENDPOINT")
        key = os.getenv("DOCUMENTINTELLIGENCE_API_KEY")
        
        if not endpoint or not key:
            raise ValueError("Document Intelligence credentials not found in environment variables")
        
        # Create output markdown file path
        output_markdown_path = Path(path).with_suffix('.md')
        
        try:            # Import the pdftomarkdown module dynamically
            spec = importlib.util.spec_from_file_location("pdftomarkdown", 
                                                         "pdftomarkdown.py")
            pdftomarkdown = importlib.util.module_from_spec(spec)
            spec.loader.exec_module(pdftomarkdown)
              # We need to modify the analyze_documents_output_in_markdown function to use our output path
            # This is a workaround since the function hardcodes the output filename to "decatur_tx_lease.md"
            # 
            # IMPLEMENTATION NOTES:
            # - We use function monkey patching to override the behavior without modifying the original
            # - This approach allows us to:
            #   1. Keep the original pdftomarkdown.py unchanged
            #   2. Control where the output is written (proper path with correct name)
            #   3. Still use all the Document Intelligence processing logic
            #
            # The wrapper performs these steps:
            # 1. Call the original function to process the PDF
            # 2. Take its output and write to our desired output path
            # 3. Return the same markdown content for other processing
            
            # Get the original function
            original_func = pdftomarkdown.analyze_documents_output_in_markdown
            
            # Create a wrapper that will direct output to our desired path
            def wrapped_func(endpoint, key, file_path, connection_string, container_name):
                # Call the original function
                markdown_content = original_func(endpoint, key, file_path, connection_string, container_name)
                
                # Now write the content to our desired location
                with open(output_markdown_path, "w", encoding="utf-8") as file:
                    file.write(markdown_content)
                
                return markdown_content
            
            # Replace the original function temporarily
            pdftomarkdown.analyze_documents_output_in_markdown = wrapped_func
            
            # Call the function with minimal parameters needed for conversion
            # We don't need Azure Blob Storage for our case
            pdftomarkdown.analyze_documents_output_in_markdown(
                endpoint=endpoint,                key=key,
                file_path=path,
                connection_string="",  # Not needed if we're not uploading
                container_name=""      # Not needed if we're not uploading
            )
            
            logger.info(f"Successfully converted image PDF to markdown: {output_markdown_path}")
            
            # Process the markdown file using the existing function
            # For OCRed PDFs, we also want optimized chunking parameters
            pdf_chunk_size = min(chunk_size * 1.5, 2000)  # Larger chunks, but max 2000 chars
            pdf_chunk_overlap = min(chunk_overlap * 2, 400)  # More overlap to maintain context
            
            docs = create_docs_from_markdown(
                path=str(output_markdown_path),
                embedding_model=embedding_model,
                chunk_size=int(pdf_chunk_size),
                chunk_overlap=int(pdf_chunk_overlap)
            )
            
            # Update filepath to point to the original PDF
            for doc in docs:
                doc['filepath'] = path
            
            return docs
            
        except Exception as e:
            logger.error(f"Error converting image PDF to markdown: {e}")
            raise ValueError(f"Failed to process image PDF: {e}")

def create_docs_from_word(path: str,
                         embedding_model,
                         chunk_size: int = 1000, 
                         chunk_overlap: int = 200) -> list[dict[str, any]]:
    """
    Process a Word document and generate vector embeddings for each chunk.
    
    This function extracts text from Word documents while preserving structure
    such as headings, paragraphs, and tables, then converts to markdown format
    for consistent processing.
    
    Args:
        path: Path to the Word document
        embedding_model: Model for generating embeddings
        chunk_size: Base size for text chunks
        chunk_overlap: Overlap between consecutive chunks
        
    Returns:
        List of document records with embeddings and metadata
    """
    if not HAS_OFFICE_SUPPORT:
        raise ImportError("python-docx is not installed. Install with: pip install python-docx")
    
    logger.info(f"Processing Word document: {path}")
    
    # Load the document
    doc = docx.Document(path)
    
    # Extract text with structure preservation
    all_text = []
    
    # Process the document title if available
    title = doc.core_properties.title
    if title:
        all_text.append(f"# {title}")
    else:
        # Use filename as title if no document title is set
        all_text.append(f"# {Path(path).stem}")
    
    # Extract document headings and paragraphs with structure
    for para in doc.paragraphs:
        if not para.text.strip():
            continue
            
        # Check if it's a heading
        if para.style.name.startswith('Heading'):
            # Get heading level from style name (e.g., 'Heading 1' -> 1)
            try:
                heading_level = int(para.style.name.split(' ')[-1])
                prefix = '#' * (heading_level + 1)  # +1 because we used H1 for the document title
            except (ValueError, IndexError):
                prefix = '##'  # Default to H2 if we can't determine the level
                
            all_text.append(f"{prefix} {para.text}")
        else:
            # Regular paragraph
            all_text.append(para.text)
    
    # Extract tables
    for table in doc.tables:
        table_rows = []
        header_row = []
        
        # Process header row
        for cell in table.rows[0].cells:
            header_row.append(cell.text.strip())
        
        if header_row:
            table_rows.append('| ' + ' | '.join(header_row) + ' |')
            # Add separator row
            table_rows.append('| ' + ' | '.join(['---' for _ in header_row]) + ' |')
        
        # Process data rows
        for row in table.rows[1:]:
            row_cells = [cell.text.strip() for cell in row.cells]
            if any(row_cells):  # Only add if row has content
                table_rows.append('| ' + ' | '.join(row_cells) + ' |')
        
        if table_rows:
            all_text.append("\n".join(table_rows))
    
    # Convert to markdown format for processing
    markdown_content = "\n\n".join(all_text)
    
    # Create a temporary markdown file
    temp_markdown_path = Path(tempfile.gettempdir()) / f"{Path(path).stem}_converted.md"
    with open(temp_markdown_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    
    logger.info(f"Created temporary markdown from Word document at: {temp_markdown_path}")
    
    # Use optimized chunking parameters for Word docs
    word_chunk_size = min(chunk_size * 1.2, 1800)  # Larger chunks for Word docs
    word_chunk_overlap = min(chunk_overlap * 1.5, 350)  # More overlap to maintain context
    
    # Use existing markdown processing
    docs = create_docs_from_markdown(
        path=str(temp_markdown_path),
        embedding_model=embedding_model,
        chunk_size=int(word_chunk_size),
        chunk_overlap=int(word_chunk_overlap)
    )
    
    # Update filepath to point to original Word document
    for doc in docs:
        doc['filepath'] = path
        doc['content_type'] = 'word'
    
    return docs

def create_docs_from_powerpoint(path: str,
                         embedding_model,
                         chunk_size: int = 1000, 
                         chunk_overlap: int = 200) -> list[dict[str, any]]:
    """
    Process a PowerPoint presentation and generate vector embeddings for each chunk.
    
    This function extracts text from PowerPoint slides while preserving structure
    such as titles, content, and notes, then converts to markdown format for 
    consistent processing.
    
    Args:
        path: Path to the PowerPoint document
        embedding_model: Model for generating embeddings
        chunk_size: Base size for text chunks
        chunk_overlap: Overlap between consecutive chunks
        
    Returns:
        List of document records with embeddings and metadata
    """
    if not HAS_OFFICE_SUPPORT:
        raise ImportError("python-pptx is not installed. Install with: pip install python-pptx")
    
    logger.info(f"Processing PowerPoint presentation: {path}")
    
    # Load the presentation
    presentation = pptx.Presentation(path)
    
    # Extract text with structure preservation
    all_text = []
    
    # Add the presentation title
    title = Path(path).stem
    all_text.append(f"# {title}")
    
    # Process each slide
    for i, slide in enumerate(presentation.slides):
        # Add slide number as heading
        slide_number = i + 1
        all_text.append(f"## Slide {slide_number}")
        
        # Extract slide title if available
        if slide.shapes.title and slide.shapes.title.text:
            all_text.append(f"### {slide.shapes.title.text.strip()}")
        
        # Extract text from all shapes in the slide
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        if texts:
            all_text.append("\n".join(texts))
        
        # Extract notes if available
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                all_text.append(f"**Notes:**\n\n{notes_text}")
    
    # Convert to markdown format for processing
    markdown_content = "\n\n".join(all_text)
    
    # Create a temporary markdown file
    temp_markdown_path = Path(tempfile.gettempdir()) / f"{Path(path).stem}_converted.md"
    with open(temp_markdown_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    
    logger.info(f"Created temporary markdown from PowerPoint at: {temp_markdown_path}")
    
    # Use optimized chunking parameters for PowerPoint docs
    # PowerPoint slides are more fragmented, so we use smaller chunks with more overlap
    ppt_chunk_size = min(chunk_size * 0.8, 1500)  # Smaller chunks than Word docs
    ppt_chunk_overlap = min(chunk_overlap * 1.5, 350)  # More overlap to maintain context
    
    # Use existing markdown processing
    docs = create_docs_from_markdown(
        path=str(temp_markdown_path),
        embedding_model=embedding_model,
        chunk_size=int(ppt_chunk_size),
        chunk_overlap=int(ppt_chunk_overlap)
    )
    
    # Update filepath to point to original PowerPoint document
    for doc in docs:
        doc['filepath'] = path
        doc['content_type'] = 'powerpoint'
    
    return docs

def create_index_from_file(index_name, file_path, file_type='markdown'):
    """
    Create an Azure AI Search index from a file.
    
    Args:
        index_name: The name of the search index to create
        file_path: Path to the file to process (markdown or csv)
        file_type: Type of file to process ('markdown' or 'csv')
    """
    # If a search index already exists, delete it:
    try:
        index_definition = index_client.get_index(index_name)
        index_client.delete_index(index_name)
        logger.info(f"🗑️  Found existing index named '{index_name}', and deleted it")
    except Exception:
        logger.info(f"🗑️  No existing index named '{index_name}' found, so no need to delete it")

    # Get the model name for vector dimensions
    model_name = os.getenv("EMBEDDINGS_MODEL", "text-embedding-ada-002")
    
    # Create an empty search index
    index_definition = create_index_definition(index_name, model=model_name)
    index_client.create_index(index_definition)
    
    # Based on file type, process the file differently
    if file_type.lower() == 'csv':
        # Create documents from the CSV file, generating vector embeddings
        docs = create_docs_from_csv(path=file_path, content_column="description", embedding_model=embeddings)
    elif file_type.lower() == 'markdown':
        # Create documents from the markdown file, chunking and generating vector embeddings
        docs = create_docs_from_markdown(path=file_path, embedding_model=embeddings)
    elif file_type.lower() == 'pdf':
        # Create documents from the PDF file, extracting text and generating vector embeddings
        docs = create_docs_from_pdf(path=file_path, embedding_model=embeddings)
    elif file_type.lower() == 'word':
        # Create documents from the Word file, extracting text and generating vector embeddings
        docs = create_docs_from_word(path=file_path, embedding_model=embeddings)
    elif file_type.lower() == 'powerpoint':
        # Create documents from the PowerPoint file, extracting text and generating vector embeddings
        docs = create_docs_from_powerpoint(path=file_path, embedding_model=embeddings)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    # Add the documents to the index using the Azure AI Search client
    search_client = SearchClient(
        endpoint=search_service_endpoint,
        index_name=index_name,
        credential=AzureKeyCredential(search_api_key),
    )

    search_client.upload_documents(docs)
    logger.info(f"➕ Uploaded {len(docs)} documents to '{index_name}' index")


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--index-name",
        type=str,
        help="index name to use when creating the AI Search index",
        default="document-index",
    )
    parser.add_argument(
        "--file", type=str, help="path to data file for creating search index", required=True
    )
    parser.add_argument(
        "--file-type", type=str, choices=["csv", "markdown", "pdf", "word", "powerpoint", "auto"], default="auto",
        help="type of file to process (csv, markdown, pdf, word, powerpoint, or auto-detect)"    )

    args = parser.parse_args()
    index_name = args.index_name
    file_path = args.file
    file_type = args.file_type
    
    # If file type is 'auto', try to detect it automatically
    if file_type.lower() == 'auto':
        is_valid, error_msg, detected_type = validate_file(file_path)
        if not is_valid:
            logger.error(f"File validation error: {error_msg}")
            exit(1)
        
        if detected_type == 'unknown':
            logger.error(f"Could not determine file type for {file_path}. Please specify --file-type explicitly.")
            exit(1)
            
        file_type = detected_type
        logger.info(f"Auto-detected file type: {file_type}")
    
    logger.info(f"Creating index from {file_type} file: {file_path}")
    
    # Additional logging for PDF files
    if file_type.lower() == 'pdf':
        logger.info("PDF processing will detect whether the file contains extractable text (native PDF)")
        logger.info("or requires OCR processing via Document Intelligence (image-based PDF)")
        logger.info("For image-based PDFs, ensure DOCUMENTINTELLIGENCE_ENDPOINT and DOCUMENTINTELLIGENCE_API_KEY are set")
    
    # Additional logging for Office documents
    if file_type.lower() in ['word', 'powerpoint'] and not HAS_OFFICE_SUPPORT:
        logger.warning(f"Processing {file_type} files requires additional libraries.")
        logger.warning("Please install the missing dependencies with: pip install python-docx python-pptx")
        exit(1)
    
    create_index_from_file(index_name, file_path, file_type)
    logger.info("Index created successfully!")